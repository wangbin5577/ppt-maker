"""
md2ppt - 从Markdown提示词自动生成LR模板PPT
============================================
支持两种格式：
  1. 标准格式: ## Slide 1 [cover]
  2. 中文格式: ## 第1页 封面

使用方法：
    python md2ppt.py <提示词.md> [输出文件.pptx]

示例：
    python md2ppt.py my_presentation.md
    python md2ppt.py my_presentation.md output.pptx

依赖：
    pip install python-pptx
"""

import sys
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt


class SlideData:
    """一页幻灯片的解析数据"""
    def __init__(self, index, layout="content"):
        self.index = index
        self.layout = layout
        self.title = ""
        self.subtitle = ""
        self.bullets = []
        self.left_title = ""
        self.left_items = []
        self.right_title = ""
        self.right_items = []
        self.table_headers = []
        self.table_rows = []
        self.notes = []
        self.code_blocks = []    # [str, ...] 代码块列表


class MarkdownParser:
    """解析PPT提示词Markdown文件，支持标准格式和中文自由格式"""

    LAYOUT_PATTERN = re.compile(r'##\s+Slide\s+(\d+)\s*\[(\w+)\]', re.IGNORECASE)
    LAYOUT_PATTERN_NO_TYPE = re.compile(r'##\s+Slide\s+(\d+)\s*$', re.IGNORECASE)
    # 英文描述格式: ## Slide 1 — Cover 或 ## Slide 1 - About Me
    LAYOUT_PATTERN_DESC = re.compile(r'##\s+Slide\s+(\d+)\s*[—\-–]+\s*(.*)', re.IGNORECASE)
    CN_PATTERN = re.compile(r'##\s+第(\d+)页\s*(.*)', re.IGNORECASE)

    CN_LAYOUT_KEYWORDS = {
        "封面": "cover",
        "结束": "end",
        "谢谢": "end",
        "章节": "section",
        "对比": "comparison",
        "比较": "comparison",
    }

    # 英文描述到布局的映射
    EN_LAYOUT_KEYWORDS = {
        "cover": "cover",
        "closing": "end",
        "thank": "end",
        "summary": "content",
    }

    def __init__(self, md_path):
        if not os.path.exists(md_path):
            raise FileNotFoundError(f"文件不存在: {md_path}")
        with open(md_path, 'r', encoding='utf-8') as f:
            self.content = f.read()
        self.slides = []

    def parse(self):
        """解析Markdown文件，返回SlideData列表"""
        self.slides = self._parse_standard()
        if self.slides:
            return self.slides
        self.slides = self._parse_chinese()
        return self.slides

    def _parse_standard(self):
        """解析标准格式 (## Slide N [layout] 或 ## Slide N — Description)"""
        slides = []
        sections = re.split(r'(?=^## Slide\s+\d+)', self.content, flags=re.MULTILINE)

        for section in sections:
            section = section.strip()
            if not section:
                continue

            # 格式1: ## Slide 1 [cover]
            match = self.LAYOUT_PATTERN.match(section)
            if match:
                index = int(match.group(1))
                layout = match.group(2).lower()
                body = section[match.end():].strip()
                slide = SlideData(index, layout)
                self._parse_slide_body(slide, body)
                slides.append(slide)
                continue

            # 格式2: ## Slide 1 — Cover
            match = self.LAYOUT_PATTERN_DESC.match(section)
            if match:
                index = int(match.group(1))
                description = match.group(2).strip()
                body = section[match.end():].strip()
                layout = self._infer_layout_from_en(description, body, index)
                slide = SlideData(index, layout)
                self._parse_slide_body(slide, body)
                slides.append(slide)
                continue

            # 格式3: ## Slide 1 (无标记)
            match = self.LAYOUT_PATTERN_NO_TYPE.match(section)
            if match:
                index = int(match.group(1))
                body = section[match.end():].strip()
                slide = SlideData(index, "content")
                self._parse_slide_body(slide, body)
                slides.append(slide)
                continue

        return slides

    def _infer_layout_from_en(self, description, body, index):
        """从英文描述推断布局类型"""
        desc_lower = description.lower()
        for keyword, layout in self.EN_LAYOUT_KEYWORDS.items():
            if keyword in desc_lower:
                return layout
        if index == 1:
            return "cover"
        # 检查是否有表格
        table_lines = [l for l in body.split('\n') if l.strip().startswith('|')]
        if len(table_lines) >= 3:
            return "table"
        # 检查是否有两列对比
        if "two-column" in body.lower() or "comparison" in desc_lower:
            return "comparison"
        return "content"

    def _parse_chinese(self):
        """解析中文格式 (## 第N页 描述)"""
        slides = []
        sections = re.split(r'(?=^## 第\d+页)', self.content, flags=re.MULTILINE)

        for section in sections:
            section = section.strip()
            if not section:
                continue

            match = self.CN_PATTERN.match(section)
            if not match:
                continue

            index = int(match.group(1))
            description = match.group(2).strip()
            body = section[match.end():].strip()

            layout = self._infer_layout_from_cn(description, body, index)
            slide = SlideData(index, layout)
            self._parse_cn_body(slide, body, description)
            slides.append(slide)

        return slides

    def _infer_layout_from_cn(self, description, body, index):
        """从中文描述和内容推断布局类型"""
        for keyword, layout in self.CN_LAYOUT_KEYWORDS.items():
            if keyword in description:
                return layout

        if index == 1:
            return "cover"

        if "对比" in body or "左右两列" in body or "两列对比" in body:
            return "comparison"

        table_lines = [l for l in body.split('\n') if l.strip().startswith('|')]
        if len(table_lines) >= 3:
            return "table"

        return "content"

    def _parse_cn_body(self, slide, body, description):
        """解析中文格式的body内容"""
        # 先提取代码块
        body, code_blocks = self._extract_code_blocks(body)
        slide.code_blocks = code_blocks

        lines = body.split('\n')
        current_section = None
        table_started = False

        for line in lines:
            stripped = line.strip()

            if not stripped or stripped == '---':
                continue

            # 标题
            if stripped.startswith('标题：') or stripped.startswith('标题:'):
                slide.title = stripped.split('：', 1)[-1].split(':', 1)[-1].strip()
                current_section = None
                continue

            if stripped.startswith('title:'):
                slide.title = stripped[6:].strip()
                current_section = None
                continue

            # 副标题
            if stripped.startswith('副标题：') or stripped.startswith('副标题:'):
                slide.subtitle = stripped.split('：', 1)[-1].split(':', 1)[-1].strip().replace('\\n', '\n')
                current_section = None
                continue

            if stripped.startswith('subtitle:'):
                slide.subtitle = stripped[9:].strip().replace('\\n', '\n')
                current_section = None
                continue

            # 要点类标记 → 开始bullet区域
            if re.match(r'^(要点|核心理念|关键|AI的作用|设计思路|工作流程|解决方案|经验教训|适用场景|可视化面板|关键数字|关键结果|每个阶段|阅读的论文|核心结论表|验证数据表|三方交叉验证|AI使用心得).*[：:]', stripped):
                current_section = "bullets"
                after_colon = stripped.split('：', 1)[-1].split(':', 1)[-1].strip()
                if after_colon:
                    slide.bullets.append((after_colon, 0))
                continue

            # 结论
            if stripped.startswith('结论：') or stripped.startswith('结论:'):
                text = stripped.split('：', 1)[-1].split(':', 1)[-1].strip()
                slide.notes.append(text)
                continue

            # 跳过配图/背景图/布局等描述
            if re.match(r'^(配图|背景图|布局|时间轴)[：:]', stripped):
                current_section = None
                continue

            # 标准格式兼容
            if stripped.startswith('left_title:'):
                slide.left_title = stripped[11:].strip()
                current_section = None
                continue
            if stripped.startswith('right_title:'):
                slide.right_title = stripped[12:].strip()
                current_section = None
                continue
            if stripped == 'left:':
                current_section = "left"
                continue
            if stripped == 'right:':
                current_section = "right"
                continue
            if stripped == 'table:':
                current_section = "table"
                table_started = False
                continue
            if stripped == 'notes:':
                current_section = "notes"
                continue

            # Markdown表格行
            if stripped.startswith('|'):
                cells = [c.strip() for c in stripped.split('|')[1:-1]]
                if not cells:
                    continue
                if all(re.match(r'^:?-+:?$', c) for c in cells):
                    table_started = True
                    continue
                if not slide.table_headers:
                    slide.table_headers = cells
                    current_section = "table"
                    table_started = False
                else:
                    slide.table_rows.append(cells)
                continue

            # bullet列表
            if stripped.startswith('- ') or stripped.startswith('* '):
                indent = len(line) - len(line.lstrip())
                level = indent // 2
                text = stripped[2:]

                if current_section == "left":
                    slide.left_items.append(text)
                elif current_section == "right":
                    slide.right_items.append(text)
                elif current_section == "notes":
                    slide.notes.append(text)
                else:
                    current_section = "bullets"
                    slide.bullets.append((text, level))
                continue

            # 有序列表
            num_match = re.match(r'^\d+[\.\)]\s*(.+)', stripped)
            if num_match:
                slide.bullets.append((num_match.group(1), 0))
                current_section = "bullets"
                continue

            # 其他文本
            if current_section == "bullets" and slide.bullets:
                last_text, last_level = slide.bullets[-1]
                slide.bullets[-1] = (last_text + " " + stripped, last_level)
            elif current_section == "left":
                slide.left_items.append(stripped)
            elif current_section == "right":
                slide.right_items.append(stripped)

        # 如果没有标题，用description
        if not slide.title and description:
            slide.title = description

        # 自动判断布局
        if slide.table_headers and slide.table_rows and slide.layout == "content":
            slide.layout = "table"

        # 如果是comparison但没有左右栏数据，降级为content
        if slide.layout == "comparison" and not slide.left_items and not slide.right_items:
            slide.layout = "content"

        # 封面页处理
        if slide.layout == "cover" and not slide.subtitle and slide.bullets:
            slide.subtitle = "\n".join([b[0] for b in slide.bullets[:2]])
            slide.bullets = []

    def _extract_code_blocks(self, body):
        """从body中提取代码块，返回 (清理后的body, [代码块列表])"""
        code_blocks = []
        # 匹配 code: 标记后跟 ``` 代码块，或者直接的 ``` 代码块
        # 模式1: code:\n```lang\n...\n```
        # 模式2: 直接 ```lang\n...\n```
        pattern = re.compile(r'(?:code:\s*\n)?```\w*\n(.*?)```', re.DOTALL)
        
        for match in pattern.finditer(body):
            code = match.group(1).rstrip('\n')
            code_blocks.append(code)
        
        # 从body中移除代码块
        cleaned = pattern.sub('', body)
        # 也移除孤立的 "code:" 行
        cleaned = re.sub(r'^code:\s*$', '', cleaned, flags=re.MULTILINE)
        
        return cleaned, code_blocks

    def _parse_slide_body(self, slide, body):
        """解析标准格式的body内容"""
        # 先提取代码块
        body, code_blocks = self._extract_code_blocks(body)
        slide.code_blocks = code_blocks

        lines = body.split('\n')
        current_section = None
        table_started = False

        for line in lines:
            stripped = line.strip()

            if not stripped:
                continue

            # title: 或 Title: 字段
            if stripped.lower().startswith('title:'):
                slide.title = stripped[6:].strip()
                current_section = None
                continue

            # subtitle: 或 Subtitle: 字段
            if stripped.lower().startswith('subtitle:'):
                slide.subtitle = stripped[9:].strip().replace('\\n', '\n')
                current_section = None
                continue

            # Key points: / Core philosophy: / How AI helped: 等 → 开始bullet
            if re.match(r'^(Key points|Core philosophy|How AI helped|Key findings|Workflow|Dashboard panels|Three AI lessons|Recommended use cases|Timeline|Layout)[:\s]', stripped, re.IGNORECASE):
                current_section = "bullets"
                continue

            # Visual: / Background image: → 跳过
            if re.match(r'^(Visual|Background image|Comparison table)[:\s]', stripped, re.IGNORECASE):
                current_section = None
                continue

            if stripped.startswith('left_title:'):
                slide.left_title = stripped[11:].strip()
                current_section = None
                continue

            if stripped.startswith('right_title:'):
                slide.right_title = stripped[12:].strip()
                current_section = None
                continue

            if stripped == 'left:':
                current_section = "left"
                continue

            if stripped == 'right:':
                current_section = "right"
                continue

            if stripped == 'table:':
                current_section = "table"
                table_started = False
                continue

            if stripped == 'notes:':
                current_section = "notes"
                continue

            if current_section == "table" and stripped.startswith('|'):
                cells = [c.strip() for c in stripped.split('|')[1:-1]]
                if all(re.match(r'^:?-+:?$', c) for c in cells):
                    table_started = True
                    continue
                if not table_started:
                    slide.table_headers = cells
                else:
                    slide.table_rows.append(cells)
                continue

            if stripped.startswith('- ') or stripped.startswith('* '):
                indent = len(line) - len(line.lstrip())
                level = indent // 2
                text = stripped[2:]

                if current_section == "left":
                    slide.left_items.append(text)
                elif current_section == "right":
                    slide.right_items.append(text)
                elif current_section == "notes":
                    slide.notes.append(text)
                else:
                    current_section = "bullets"
                    slide.bullets.append((text, level))
                continue

            # 有序列表 (1. xxx)
            num_match = re.match(r'^\d+[\.\)]\s*(.+)', stripped)
            if num_match:
                current_section = "bullets"
                slide.bullets.append((num_match.group(1), 0))
                continue

            # 跳过分隔线
            if stripped == '---':
                continue

            if current_section == "bullets" and slide.bullets:
                last_text, last_level = slide.bullets[-1]
                slide.bullets[-1] = (last_text + " " + stripped, last_level)
            elif current_section == "left":
                slide.left_items.append(stripped)
            elif current_section == "right":
                slide.right_items.append(stripped)

        # 后处理：如果有表格但布局是content，升级为table
        if slide.table_headers and slide.table_rows and slide.layout == "content":
            slide.layout = "table"
        # 如果是comparison但没有左右栏数据，降级为content
        if slide.layout == "comparison" and not slide.left_items and not slide.right_items:
            slide.layout = "content"


class PPTGenerator:
    """基于模板生成PPT，支持自动检测布局"""

    # LR模板的默认布局索引
    DEFAULT_LAYOUTS = {
        "cover": 7,
        "content": 9,
        "content_spacious": 10,
        "comparison": 15,
        "section": 14,
        "end": 7,
        "table": 10,
    }

    def __init__(self, template_path="LR模板.pptx"):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        self.prs = Presentation(template_path)
        self.LAYOUTS = self._detect_layouts()
        self._clear_existing_slides()

    def _detect_layouts(self):
        """自动检测模板中的布局类型"""
        layouts = {}
        n_layouts = len(self.prs.slide_layouts)

        # 按名称匹配
        for i, layout in enumerate(self.prs.slide_layouts):
            name = layout.name.lower()
            placeholders = list(layout.placeholders)
            has_center_title = any(ph.placeholder_format.type == 3 for ph in placeholders)
            has_title = any(ph.placeholder_format.type == 1 for ph in placeholders)
            has_object = any(ph.placeholder_format.type == 7 for ph in placeholders)
            has_body = any(ph.placeholder_format.type == 2 for ph in placeholders)
            n_ph = len(placeholders)

            # 封面：有CENTER_TITLE的布局
            if has_center_title:
                layouts["cover"] = i  # 取最后一个（LR模板中Marine GTC在后面）

            # 正文：TITLE + OBJECT，2个占位符
            if has_title and has_object and n_ph == 2:
                if "content" not in layouts:
                    layouts["content"] = i
                elif "content_spacious" not in layouts:
                    layouts["content_spacious"] = i

            # 双栏：5个以上占位符
            if n_ph >= 5 and has_title and "comparison" not in layouts:
                layouts["comparison"] = i

            # 章节：TITLE + BODY，2个占位符，无OBJECT
            if has_title and has_body and n_ph == 2 and not has_object:
                if "section" not in layouts:
                    layouts["section"] = i

        # 按名称优先匹配（覆盖上面的通用检测）
        for i, layout in enumerate(self.prs.slide_layouts):
            name = layout.name.lower()
            if "gtc" in name or "new cover" in name:
                layouts["cover"] = i
            elif "spacious" in name and "content_spacious" not in layouts:
                layouts["content_spacious"] = i
            elif "comparison" in name:
                layouts["comparison"] = i
            elif "section" in name:
                layouts["section"] = i

        # 填充缺失的
        if "cover" not in layouts:
            layouts["cover"] = 0
        if "content" not in layouts:
            layouts["content"] = min(1, n_layouts - 1)
        if "content_spacious" not in layouts:
            layouts["content_spacious"] = layouts["content"]
        if "comparison" not in layouts:
            layouts["comparison"] = layouts["content"]
        if "section" not in layouts:
            layouts["section"] = layouts["cover"]
        if "end" not in layouts:
            layouts["end"] = layouts["cover"]
        if "table" not in layouts:
            layouts["table"] = layouts["content_spacious"]

        return layouts

    def _clear_existing_slides(self):
        """清除模板中已有的幻灯片"""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    def generate(self, slides_data):
        """根据解析数据生成所有幻灯片"""
        for sd in slides_data:
            if sd.layout in ("cover", "end"):
                self._add_cover(sd)
            elif sd.layout == "comparison":
                self._add_comparison(sd)
            elif sd.layout == "table":
                self._add_table(sd)
            elif sd.layout == "section":
                self._add_section(sd)
            elif sd.layout == "content_spacious":
                self._add_content(sd, spacious=True)
            else:
                self._add_content(sd, spacious=False)

    def _add_title_background(self, slide):
        """为缺少标题背景的布局添加深蓝色标题栏和蓝色横线"""
        from pptx.util import Inches, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        
        # 深蓝色标题背景 (accent2 = #003C71)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(12), Inches(1.28)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0x00, 0x3C, 0x71)
        bg_shape.line.fill.background()
        # 移到最底层
        slide.shapes._spTree.remove(bg_shape._element)
        slide.shapes._spTree.insert(2, bg_shape._element)
        
        # 蓝色横线 (accent1 = #3B8EDE)
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.28),
            Inches(12), Inches(0.19)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(0x3B, 0x8E, 0xDE)
        line_shape.line.fill.background()
        slide.shapes._spTree.remove(line_shape._element)
        slide.shapes._spTree.insert(3, line_shape._element)

    def _add_cover(self, sd):
        layout = self.prs.slide_layouts[self.LAYOUTS["cover"]]
        slide = self.prs.slides.add_slide(layout)
        slide.placeholders[0].text = sd.title
        slide.placeholders[1].text = sd.subtitle

    def _calc_font_size(self, bullets, has_code=False):
        """根据内容量自动计算字体大小和行距"""
        n_bullets = len(bullets)
        total_chars = sum(len(text) for text, _ in bullets)

        if has_code:
            # 有代码块时，正文区域更小，字体适当缩小
            if n_bullets <= 4:
                return Pt(18), Pt(16), Pt(6), Pt(3)
            else:
                return Pt(16), Pt(14), Pt(4), Pt(2)

        if n_bullets <= 6 and total_chars < 350:
            return Pt(20), Pt(18), Pt(8), Pt(4)
        elif n_bullets <= 10 and total_chars < 600:
            return Pt(18), Pt(16), Pt(6), Pt(3)
        elif n_bullets <= 15 and total_chars < 900:
            return Pt(16), Pt(14), Pt(4), Pt(2)
        else:
            return Pt(14), Pt(13), Pt(3), Pt(1)

    def _add_content(self, sd, spacious=False):
        # 有代码块时强制用大内容区
        if sd.code_blocks:
            spacious = True
        layout_idx = self.LAYOUTS["content_spacious"] if spacious else self.LAYOUTS["content"]
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)
        slide.placeholders[0].text = sd.title

        # 自动计算字体大小
        size_l0, size_l1, space_after, space_before = self._calc_font_size(
            sd.bullets, has_code=bool(sd.code_blocks)
        )

        tf = slide.placeholders[1].text_frame
        tf.clear()

        first_para = True
        for text, level in sd.bullets:
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.size = size_l0 if level == 0 else size_l1
            p.space_after = space_after
            p.space_before = space_before

        # 代码块用独立文本框，位置基于bullet数量估算
        if sd.code_blocks:
            self._add_code_blocks(slide, sd.code_blocks, n_bullets=len(sd.bullets),
                                  font_size=size_l0)

    def _add_code_blocks(self, slide, code_blocks, n_bullets=0, font_size=Pt(18)):
        """在幻灯片上添加代码块文本框（带浅灰背景+边框）"""
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        # 估算正文占用高度：每个bullet约 (字号/72 + 段间距) 英寸
        font_inches = font_size.inches if hasattr(font_size, 'inches') else font_size / 914400
        line_height = font_inches + 0.15  # 字高 + 间距
        text_height = n_bullets * line_height
        
        # 代码块起始位置 = 内容区顶部(1.7) + 正文高度 + 间隔
        top = Inches(1.7) + Inches(text_height) + Inches(0.3)
        top = max(top, Inches(2.2))  # 最少从2.2开始
        
        page_bottom = Inches(6.6)

        for code in code_blocks:
            lines = code.split('\n')
            n_lines = len(lines)
            
            # 代码块高度：每行0.2英寸 + 上下padding
            height = Inches(n_lines * 0.2 + 0.25)
            
            # 限制不超出页面
            available = page_bottom - top
            if available < Inches(0.8):
                break
            height = min(height, available)

            # 创建文本框
            txBox = slide.shapes.add_textbox(
                Inches(0.8), top,
                Inches(10.4), height
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            # 内边距
            tf.margin_top = Pt(6)
            tf.margin_bottom = Pt(6)
            tf.margin_left = Pt(8)

            # 浅灰背景
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

            # 边框
            txBox.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            txBox.line.width = Pt(0.75)

            # 代码字号根据行数调整
            if n_lines <= 8:
                code_size = Pt(11)
            elif n_lines <= 12:
                code_size = Pt(10)
            else:
                code_size = Pt(9)

            # 填入代码
            for i, code_line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = code_line
                p.font.name = 'Consolas'
                p.font.size = code_size
                p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
                p.space_after = Pt(1)
                p.space_before = Pt(0)

            top += height + Inches(0.15)

    def _add_comparison(self, sd):
        layout = self.prs.slide_layouts[self.LAYOUTS["comparison"]]
        slide = self.prs.slides.add_slide(layout)
        
        # 添加标题背景色块（Comparison布局缺少这些装饰元素）
        self._add_title_background(slide)
        
        slide.placeholders[0].text = sd.title
        slide.placeholders[1].text = sd.left_title
        slide.placeholders[3].text = sd.right_title

        tf_left = slide.placeholders[2].text_frame
        tf_left.clear()
        for i, item in enumerate(sd.left_items):
            p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_after = Pt(8)

        tf_right = slide.placeholders[4].text_frame
        tf_right.clear()
        for i, item in enumerate(sd.right_items):
            p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_after = Pt(8)

    def _add_table(self, sd):
        layout = self.prs.slide_layouts[self.LAYOUTS["table"]]
        slide = self.prs.slides.add_slide(layout)
        slide.placeholders[0].text = sd.title

        if sd.table_headers and sd.table_rows:
            n_cols = len(sd.table_headers)
            n_rows = len(sd.table_rows) + 1

            left = Inches(0.8)
            top = Inches(1.8)
            width = Inches(10.4)
            row_height = 0.45
            height = Inches(min(row_height * n_rows, 4.5))

            table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
            table = table_shape.table

            for i, header in enumerate(sd.table_headers):
                cell = table.cell(0, i)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)

            for row_idx, row_data in enumerate(sd.table_rows, 1):
                for col_idx in range(min(len(row_data), n_cols)):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = row_data[col_idx]
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)

        if sd.notes:
            n_rows_total = (len(sd.table_rows) + 1) if sd.table_headers else 0
            note_top = Inches(1.8 + 0.45 * n_rows_total + 0.3)
            txBox = slide.shapes.add_textbox(Inches(0.8), note_top, Inches(10.4), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, note in enumerate(sd.notes):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = note

    def _add_section(self, sd):
        layout = self.prs.slide_layouts[self.LAYOUTS["section"]]
        slide = self.prs.slides.add_slide(layout)
        slide.placeholders[0].text = sd.title
        if sd.subtitle:
            slide.placeholders[1].text = sd.subtitle

    def save(self, output_path):
        self.prs.save(output_path)
        print(f"✅ PPT已生成: {output_path}")
        print(f"   共 {len(self.prs.slides)} 页")


def main():
    if len(sys.argv) < 2 or sys.argv[1] in ('-h', '--help'):
        print("用法: python md2ppt.py <提示词.md> [输出文件.pptx] [--template 模板.pptx]")
        print("示例:")
        print("  python md2ppt.py my_presentation.md")
        print("  python md2ppt.py my_presentation.md output.pptx")
        print("  python md2ppt.py my_presentation.md output.pptx --template 其他模板.pptx")
        sys.exit(0 if sys.argv[1:] and sys.argv[1] in ('-h', '--help') else 1)

    # 解析参数
    args = [a for a in sys.argv[1:] if not a.startswith('--')]
    template_path = "LR模板.pptx"
    if '--template' in sys.argv:
        idx = sys.argv.index('--template')
        if idx + 1 < len(sys.argv):
            template_path = sys.argv[idx + 1]

    md_path = args[0]
    if len(args) >= 2:
        output_path = args[1]
    else:
        name = os.path.splitext(os.path.basename(md_path))[0]
        output_path = f"{name}_output.pptx"

    print(f"📄 提示词文件: {md_path}")
    print(f"📋 模板文件: {template_path}")
    print(f"📂 输出文件: {output_path}")
    print()

    print("🔍 解析提示词...")
    parser = MarkdownParser(md_path)
    slides_data = parser.parse()
    print(f"   解析出 {len(slides_data)} 页")
    for sd in slides_data:
        print(f"   Slide {sd.index:2d} [{sd.layout:18s}] {sd.title[:50]}")
    print()

    print("🔄 生成PPT...")
    gen = PPTGenerator(template_path)
    gen.generate(slides_data)
    gen.save(output_path)


if __name__ == "__main__":
    main()
