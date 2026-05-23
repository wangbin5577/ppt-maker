"""
PPT Maker Web UI
================
基于Streamlit的Web界面，用户上传提示词.md和PPT模板，一键生成PPT。

启动方式：
    streamlit run app.py
"""

import streamlit as st
import tempfile
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt


# ============================================================
# 解析器和生成器（从 md2ppt.py 整合）
# ============================================================

class SlideData:
    """一页幻灯片的解析数据"""
    def __init__(self, index, layout="content"):
        self.index = index
        self.layout = layout
        self.title = ""
        self.subtitle = ""
        self.bullets = []
        self.left_title = ""
        self.left_items = []
        self.right_title = ""
        self.right_items = []
        self.table_headers = []
        self.table_rows = []
        self.notes = []
        self.code_blocks = []
        self.image_placeholders = []


class MarkdownParser:
    """解析PPT提示词Markdown文件，支持标准格式和中文自由格式"""

    # 标准格式: ## Slide 1 [cover]
    LAYOUT_PATTERN = re.compile(r'##\s+Slide\s+(\d+)\s*\[(\w+)\]', re.IGNORECASE)
    LAYOUT_PATTERN_NO_TYPE = re.compile(r'##\s+Slide\s+(\d+)\s*$', re.IGNORECASE)
    # 中文格式: ## 第1页 封面  或  ## 第1页 描述文字
    CN_PATTERN = re.compile(r'##\s+第(\d+)页\s*(.*)', re.IGNORECASE)

    # 中文布局关键词映射
    CN_LAYOUT_KEYWORDS = {
        "封面": "cover",
        "结束": "end",
        "谢谢": "end",
        "章节": "section",
        "对比": "comparison",
        "比较": "comparison",
    }

    def __init__(self, md_content):
        self.content = md_content
        self.slides = []

    def parse(self):
        """解析Markdown内容，返回SlideData列表"""
        # 先尝试标准格式
        self.slides = self._parse_standard()
        if self.slides:
            return self.slides

        # 再尝试中文格式
        self.slides = self._parse_chinese()
        return self.slides

    def _parse_standard(self):
        """解析标准格式 (## Slide N [layout])"""
        slides = []
        sections = re.split(r'(?=^## Slide\s+\d+)', self.content, flags=re.MULTILINE)

        for section in sections:
            section = section.strip()
            if not section:
                continue

            match = self.LAYOUT_PATTERN.match(section)
            if match:
                index = int(match.group(1))
                layout = match.group(2).lower()
                body = section[match.end():].strip()
                slide = SlideData(index, layout)
                self._parse_slide_body(slide, body)
                slides.append(slide)
                continue

            match = self.LAYOUT_PATTERN_NO_TYPE.match(section)
            if match:
                index = int(match.group(1))
                body = section[match.end():].strip()
                slide = SlideData(index, "content")
                self._parse_slide_body(slide, body)
                slides.append(slide)
                continue

        return slides

    def _parse_chinese(self):
        """解析中文格式 (## 第N页 描述)"""
        slides = []
        sections = re.split(r'(?=^## 第\d+页)', self.content, flags=re.MULTILINE)

        for section in sections:
            section = section.strip()
            if not section:
                continue

            match = self.CN_PATTERN.match(section)
            if not match:
                continue

            index = int(match.group(1))
            description = match.group(2).strip()
            body = section[match.end():].strip()

            # 根据描述推断布局
            layout = self._infer_layout_from_cn(description, body, index)
            slide = SlideData(index, layout)
            self._parse_cn_body(slide, body, description)
            slides.append(slide)

        return slides

    def _infer_layout_from_cn(self, description, body, index):
        """从中文描述和内容推断布局类型"""
        # 检查关键词
        for keyword, layout in self.CN_LAYOUT_KEYWORDS.items():
            if keyword in description:
                return layout

        # 第1页通常是封面
        if index == 1:
            return "cover"

        # 包含"对比表"或"左右两列"的是comparison
        if "对比" in body or "左右两列" in body or "两列对比" in body:
            return "comparison"

        # 包含markdown表格的是table
        table_lines = [l for l in body.split('\n') if l.strip().startswith('|')]
        if len(table_lines) >= 3:
            return "table"

        return "content"

    def _parse_cn_body(self, slide, body, description):
        """解析中文格式的body内容"""
        lines = body.split('\n')
        current_section = None
        table_started = False

        for line in lines:
            stripped = line.strip()

            if not stripped or stripped == '---':
                continue

            # 标题：xxx 或 title: xxx
            if stripped.startswith('标题：') or stripped.startswith('标题:'):
                slide.title = stripped.split('：', 1)[-1].split(':', 1)[-1].strip()
                current_section = None
                continue

            if stripped.startswith('title:'):
                slide.title = stripped[6:].strip()
                current_section = None
                continue

            # 副标题
            if stripped.startswith('副标题：') or stripped.startswith('副标题:'):
                slide.subtitle = stripped.split('：', 1)[-1].split(':', 1)[-1].strip().replace('\\n', '\n')
                current_section = None
                continue

            if stripped.startswith('subtitle:'):
                slide.subtitle = stripped[9:].strip().replace('\\n', '\n')
                current_section = None
                continue

            # 要点：/ 核心理念：/ AI的作用：/ 设计思路：等 → 开始bullet区域
            if re.match(r'^(要点|核心理念|关键|AI的作用|设计思路|工作流程|解决方案|经验教训|适用场景|可视化面板|关键数字|关键结果|每个阶段|阅读的论文|核心结论表|验证数据表|三方交叉验证|AI使用心得).*[：:]', stripped):
                current_section = "bullets"
                # 如果冒号后面有内容，作为一个bullet
                after_colon = stripped.split('：', 1)[-1].split(':', 1)[-1].strip()
                if after_colon:
                    slide.bullets.append((after_colon, 0))
                continue

            # 对比表：开始表格
            if re.match(r'^(对比表|核心结论表|验证数据表|三方交叉验证结果)[：:]?', stripped):
                current_section = "table"
                table_started = False
                continue

            # 结论：/ 经验教训：→ 作为备注
            if stripped.startswith('结论：') or stripped.startswith('结论:'):
                text = stripped.split('：', 1)[-1].split(':', 1)[-1].strip()
                slide.notes.append(text)
                continue

            # 配图：/ 背景图：→ 跳过（程序不处理图片）
            if re.match(r'^(配图|背景图|布局|时间轴)[：:]', stripped):
                current_section = None
                continue

            # left_title / right_title (标准格式兼容)
            if stripped.startswith('left_title:'):
                slide.left_title = stripped[11:].strip()
                current_section = None
                continue
            if stripped.startswith('right_title:'):
                slide.right_title = stripped[12:].strip()
                current_section = None
                continue
            if stripped == 'left:':
                current_section = "left"
                continue
            if stripped == 'right:':
                current_section = "right"
                continue
            if stripped == 'table:':
                current_section = "table"
                table_started = False
                continue
            if stripped == 'notes:':
                current_section = "notes"
                continue

            # Markdown表格行
            if stripped.startswith('|'):
                cells = [c.strip() for c in stripped.split('|')[1:-1]]
                if not cells:
                    continue
                if all(re.match(r'^:?-+:?$', c) for c in cells):
                    table_started = True
                    continue
                if not slide.table_headers:
                    slide.table_headers = cells
                    current_section = "table"
                    table_started = False
                else:
                    slide.table_rows.append(cells)
                continue

            # bullet列表
            if stripped.startswith('- ') or stripped.startswith('* '):
                indent = len(line) - len(line.lstrip())
                level = indent // 2
                text = stripped[2:]

                if current_section == "left":
                    slide.left_items.append(text)
                elif current_section == "right":
                    slide.right_items.append(text)
                elif current_section == "notes":
                    slide.notes.append(text)
                else:
                    current_section = "bullets"
                    slide.bullets.append((text, level))
                continue

            # 有序列表 (1. xxx)
            num_match = re.match(r'^\d+[\.\)]\s*(.+)', stripped)
            if num_match:
                slide.bullets.append((num_match.group(1), 0))
                current_section = "bullets"
                continue

            # 其他文本 — 如果在bullets模式追加，否则作为新bullet
            if current_section == "bullets" and slide.bullets:
                last_text, last_level = slide.bullets[-1]
                slide.bullets[-1] = (last_text + " " + stripped, last_level)
            elif current_section == "left":
                slide.left_items.append(stripped)
            elif current_section == "right":
                slide.right_items.append(stripped)

        # 如果没有提取到标题，用description作为标题
        if not slide.title and description:
            slide.title = description

        # 自动判断：如果有表格数据但布局不是table，更新布局
        if slide.table_headers and slide.table_rows and slide.layout == "content":
            slide.layout = "table"

        # 如果是comparison但没有左右栏数据，降级为content
        if slide.layout == "comparison" and not slide.left_items and not slide.right_items:
            slide.layout = "content"

        # 封面页：如果没有副标题，把第一个bullet当副标题
        if slide.layout == "cover" and not slide.subtitle and slide.bullets:
            slide.subtitle = "\n".join([b[0] for b in slide.bullets[:2]])
            slide.bullets = []

    def _extract_code_blocks(self, body):
        """从body中提取代码块"""
        code_blocks = []
        pattern = re.compile(r'(?:code:\s*\n)?```\w*\n(.*?)```', re.DOTALL)
        for match in pattern.finditer(body):
            code = match.group(1).rstrip('\n')
            code_blocks.append(code)
        cleaned = pattern.sub('', body)
        cleaned = re.sub(r'^code:\s*$', '', cleaned, flags=re.MULTILINE)
        return cleaned, code_blocks

    def _parse_slide_body(self, slide, body):
        """解析标准格式的body内容"""
        body, code_blocks = self._extract_code_blocks(body)
        slide.code_blocks = code_blocks
        lines = body.split('\n')
        current_section = None
        table_started = False

        for line in lines:
            stripped = line.strip()

            if not stripped:
                continue

            if stripped.startswith('title:'):
                slide.title = stripped[6:].strip()
                current_section = None
                continue

            if stripped.startswith('subtitle:'):
                slide.subtitle = stripped[9:].strip().replace('\\n', '\n')
                current_section = None
                continue

            if stripped.lower().startswith('image:'):
                caption = stripped[6:].strip()
                slide.image_placeholders.append(caption)
                current_section = None
                continue

            if stripped.startswith('left_title:'):
                slide.left_title = stripped[11:].strip()
                current_section = None
                continue

            if stripped.startswith('right_title:'):
                slide.right_title = stripped[12:].strip()
                current_section = None
                continue

            if stripped == 'left:':
                current_section = "left"
                continue

            if stripped == 'right:':
                current_section = "right"
                continue

            if stripped == 'table:':
                current_section = "table"
                table_started = False
                continue

            if stripped == 'notes:':
                current_section = "notes"
                continue

            if current_section == "table" and stripped.startswith('|'):
                cells = [c.strip() for c in stripped.split('|')[1:-1]]
                if all(re.match(r'^:?-+:?$', c) for c in cells):
                    table_started = True
                    continue
                if not table_started:
                    slide.table_headers = cells
                else:
                    slide.table_rows.append(cells)
                continue

            if stripped.startswith('- ') or stripped.startswith('* '):
                indent = len(line) - len(line.lstrip())
                level = indent // 2
                text = stripped[2:]

                if current_section == "left":
                    slide.left_items.append(text)
                elif current_section == "right":
                    slide.right_items.append(text)
                elif current_section == "notes":
                    slide.notes.append(text)
                else:
                    current_section = "bullets"
                    slide.bullets.append((text, level))
                continue

            if current_section == "bullets" and slide.bullets:
                last_text, last_level = slide.bullets[-1]
                slide.bullets[-1] = (last_text + " " + stripped, last_level)
            elif current_section == "left":
                slide.left_items.append(stripped)
            elif current_section == "right":
                slide.right_items.append(stripped)


class PPTGenerator:
    """基于模板生成PPT"""

    # 默认LR模板布局索引（可通过自动检测覆盖）
    DEFAULT_LAYOUTS = {
        "cover": 7,
        "content": 9,
        "content_spacious": 10,
        "comparison": 15,
        "section": 14,
        "end": 7,
        "table": 10,
    }

    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.layouts = self._detect_layouts()
        self._clear_existing_slides()

    def _detect_layouts(self):
        """自动检测模板中的布局类型，优先使用布局名称匹配"""
        layouts = {}
        
        # 第一轮：按名称精确匹配（适用于LR模板等已知模板）
        name_mapping = {
            "marine gtc": "cover",
            "marine turbine": "cover",
            "marine compass": "cover",
            "marine propeller": "cover",
            "marine newbuilding": "cover",
            "marine engine room": "cover",
            "marine safety": "cover",
            "marine waves": "cover",
            "new cover": "cover",
            "white standard": "content",
            "white spacious": "content_spacious",
            "grey standard": "content",
            "grey spacous": "content_spacious",
            "comparison": "comparison",
            "section header": "section",
            "end slide": "end",
            "title only": "title_only",
            "blank": "blank",
        }
        
        for i, layout in enumerate(self.prs.slide_layouts):
            name_lower = layout.name.lower()
            for pattern, layout_type in name_mapping.items():
                if pattern in name_lower:
                    # 只保留第一个匹配的（除了cover，保留最后一个以匹配Marine GTC）
                    if layout_type == "cover":
                        layouts["cover"] = i  # 会被后面的覆盖，最终取最后一个cover布局
                    elif layout_type not in layouts:
                        layouts[layout_type] = i
                    break
        
        # 对于cover，优先选择 "Marine GTC" 或 "New Cover"
        for i, layout in enumerate(self.prs.slide_layouts):
            name_lower = layout.name.lower()
            if "gtc" in name_lower or "new cover" in name_lower:
                layouts["cover"] = i
                break
        
        # 第二轮：如果名称匹配不够，用占位符类型检测补充
        if "content" not in layouts or "comparison" not in layouts:
            for i, layout in enumerate(self.prs.slide_layouts):
                placeholders = list(layout.placeholders)
                has_center_title = any(ph.placeholder_format.type == 3 for ph in placeholders)
                has_title = any(ph.placeholder_format.type == 1 for ph in placeholders)
                has_object = any(ph.placeholder_format.type == 7 for ph in placeholders)
                has_body = any(ph.placeholder_format.type == 2 for ph in placeholders)
                n_placeholders = len(placeholders)

                if has_title and has_object and n_placeholders == 2:
                    if "content" not in layouts:
                        layouts["content"] = i
                    elif "content_spacious" not in layouts:
                        layouts["content_spacious"] = i
                if n_placeholders >= 5 and has_title and "comparison" not in layouts:
                    layouts["comparison"] = i
                if has_title and has_body and n_placeholders == 2 and not has_object:
                    if "section" not in layouts:
                        layouts["section"] = i

        # 用默认值填充缺失的
        for key, val in self.DEFAULT_LAYOUTS.items():
            if key not in layouts:
                if val < len(self.prs.slide_layouts):
                    layouts[key] = val
                else:
                    layouts[key] = layouts.get("content", 0)

        if "end" not in layouts:
            layouts["end"] = layouts.get("cover", 0)
        if "table" not in layouts:
            layouts["table"] = layouts.get("content_spacious", layouts.get("content", 0))

        return layouts

    def _clear_existing_slides(self):
        """清除模板中已有的幻灯片"""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    def generate(self, slides_data):
        """根据解析数据生成所有幻灯片"""
        for sd in slides_data:
            if sd.layout in ("cover", "end"):
                self._add_cover(sd)
            elif sd.layout == "comparison":
                self._add_comparison(sd)
            elif sd.layout == "table":
                self._add_table(sd)
            elif sd.layout == "section":
                self._add_section(sd)
            elif sd.layout == "content_spacious":
                self._add_content(sd, spacious=True)
            else:
                self._add_content(sd, spacious=False)

    def _add_title_background(self, slide):
        """为缺少标题背景的布局添加深蓝色标题栏和蓝色横线，并把标题文字设为白色"""
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
        from lxml import etree

        # 深蓝色标题背景 (accent2 = #003C71)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(12), Inches(1.28)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0x00, 0x3C, 0x71)
        bg_shape.line.fill.background()
        slide.shapes._spTree.remove(bg_shape._element)
        slide.shapes._spTree.insert(2, bg_shape._element)

        # 蓝色横线 (accent1 = #3B8EDE)
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.28),
            Inches(12), Inches(0.19)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(0x3B, 0x8E, 0xDE)
        line_shape.line.fill.background()
        slide.shapes._spTree.remove(line_shape._element)
        slide.shapes._spTree.insert(3, line_shape._element)

        # 通过XML设置标题占位符默认白色字体
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                txBody = ph._element.find(qn('p:txBody'))
                if txBody is not None:
                    lstStyle = txBody.find(qn('a:lstStyle'))
                    if lstStyle is None:
                        lstStyle = etree.SubElement(txBody, qn('a:lstStyle'))
                    defPPr = lstStyle.find(qn('a:defPPr'))
                    if defPPr is None:
                        defPPr = etree.SubElement(lstStyle, qn('a:defPPr'))
                    defRPr = defPPr.find(qn('a:defRPr'))
                    if defRPr is None:
                        defRPr = etree.SubElement(defPPr, qn('a:defRPr'))
                    solidFill = defRPr.find(qn('a:solidFill'))
                    if solidFill is None:
                        solidFill = etree.SubElement(defRPr, qn('a:solidFill'))
                    srgbClr = solidFill.find(qn('a:srgbClr'))
                    if srgbClr is None:
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                    srgbClr.set('val', 'FFFFFF')
                break

    def _set_title_white(self, slide):
        """把标题占位符的文字颜色设为白色（在设置文字之后调用）"""
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        from lxml import etree
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                tf = ph.text_frame
                for para in tf.paragraphs:
                    para.font.color.rgb = WHITE
                    for run in para.runs:
                        run.font.color.rgb = WHITE
                    for r in para._p.findall(qn('a:r')):
                        rPr = r.find(qn('a:rPr'))
                        if rPr is None:
                            rPr = etree.SubElement(r, qn('a:rPr'))
                            r.insert(0, rPr)
                        for old in rPr.findall(qn('a:solidFill')):
                            rPr.remove(old)
                        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                        srgbClr.set('val', 'FFFFFF')
                break

    def _add_cover(self, sd):
        layout = self.prs.slide_layouts[self.layouts["cover"]]
        slide = self.prs.slides.add_slide(layout)
        if 0 in [ph.placeholder_format.idx for ph in slide.placeholders]:
            slide.placeholders[0].text = sd.title
        if 1 in [ph.placeholder_format.idx for ph in slide.placeholders]:
            slide.placeholders[1].text = sd.subtitle

    def _calc_font_size(self, bullets, has_code=False):
        """根据内容量自动计算字体大小和行距"""
        n_bullets = len(bullets)
        total_chars = sum(len(text) for text, _ in bullets)

        if has_code:
            if n_bullets <= 4:
                return Pt(18), Pt(16), Pt(6), Pt(3)
            else:
                return Pt(16), Pt(14), Pt(4), Pt(2)

        if n_bullets <= 6 and total_chars < 350:
            return Pt(20), Pt(18), Pt(8), Pt(4)
        elif n_bullets <= 10 and total_chars < 600:
            return Pt(18), Pt(16), Pt(6), Pt(3)
        elif n_bullets <= 15 and total_chars < 900:
            return Pt(16), Pt(14), Pt(4), Pt(2)
        else:
            return Pt(14), Pt(13), Pt(3), Pt(1)

    def _add_content(self, sd, spacious=False):
        has_images = bool(sd.image_placeholders)
        has_code = bool(sd.code_blocks)

        if has_images and not has_code:
            self._add_content_with_images(sd)
            return

        if has_code:
            self._add_content_with_code(sd)
            return

        layout_key = "content_spacious" if spacious else "content"
        layout = self.prs.slide_layouts[self.layouts[layout_key]]
        slide = self.prs.slides.add_slide(layout)

        if 0 in [ph.placeholder_format.idx for ph in slide.placeholders]:
            slide.placeholders[0].text = sd.title

        size_l0, size_l1, space_after, space_before = self._calc_font_size(sd.bullets)

        if 1 in [ph.placeholder_format.idx for ph in slide.placeholders]:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            first_para = True
            for text, level in sd.bullets:
                if first_para:
                    p = tf.paragraphs[0]
                    first_para = False
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.level = level
                p.font.size = size_l0 if level == 0 else size_l1
                p.space_after = space_after
                p.space_before = space_before

    def _add_content_with_code(self, sd):
        """有代码块时使用左文右代码的分栏布局"""
        from pptx.dml.color import RGBColor
        layout_key = "title_only" if "title_only" in self.layouts else "content"
        layout = self.prs.slide_layouts[self.layouts[layout_key]]
        slide = self.prs.slides.add_slide(layout)

        ph_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        if 0 in ph_indices:
            slide.placeholders[0].text = sd.title

        # 添加标题背景
        self._add_title_background(slide)
        self._set_title_white(slide)

        # 删除所有内容占位符
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx != 0:
                ph._element.getparent().remove(ph._element)

        content_top = Inches(1.65)
        content_bottom = Inches(6.6)
        content_height = content_bottom - content_top
        left_margin = Inches(0.5)
        gap = Inches(0.3)
        total_width = Inches(11.2)
        n_bullets = len(sd.bullets)
        n_codes = len(sd.code_blocks)

        if n_bullets == 0:
            code_width = total_width
            code_left = left_margin
        else:
            text_width = total_width * 0.40
            code_width = total_width * 0.55
            code_left = left_margin + text_width + gap

            size_l0, size_l1, space_after, space_before = self._calc_font_size(sd.bullets)
            txBox = slide.shapes.add_textbox(left_margin, content_top, text_width, content_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_top = Pt(4)
            tf.margin_left = Pt(4)
            for i, (text, level) in enumerate(sd.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = ("  " * level + "• " if level > 0 else "• ") + text
                p.font.size = size_l0 if level == 0 else size_l1
                p.font.color.rgb = RGBColor(0x1A, 0x6B, 0xB5)
                p.space_after = space_after
                p.space_before = space_before

        if n_codes > 0:
            code_gap = Inches(0.15)
            single_height = (content_height - code_gap * (n_codes - 1)) / n_codes
            top = content_top
            for code in sd.code_blocks:
                lines = code.split('\n')
                n_lines = len(lines)
                height = min(single_height, Inches(n_lines * 0.22 + 0.3))
                height = max(height, Inches(0.8))

                txBox = slide.shapes.add_textbox(code_left, top, code_width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.margin_top = Pt(6)
                tf.margin_bottom = Pt(6)
                tf.margin_left = Pt(8)
                txBox.fill.solid()
                txBox.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
                txBox.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
                txBox.line.width = Pt(0.75)

                code_size = Pt(11) if n_lines <= 8 else Pt(10) if n_lines <= 14 else Pt(9)
                for i, code_line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = code_line
                    p.font.name = 'Consolas'
                    p.font.size = code_size
                    p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
                    p.space_after = Pt(1)
                    p.space_before = Pt(0)
                top += height + code_gap

    def _add_content_with_images(self, sd):
        """有图片时使用左文右图的分栏布局"""
        layout_key = "title_only" if "title_only" in self.layouts else "content"
        layout = self.prs.slide_layouts[self.layouts[layout_key]]
        slide = self.prs.slides.add_slide(layout)

        ph_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        if 0 in ph_indices:
            slide.placeholders[0].text = sd.title

        # 添加标题背景
        self._add_title_background(slide)
        self._set_title_white(slide)

        # 删除所有内容占位符
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx != 0:
                ph._element.getparent().remove(ph._element)

        n_images = len(sd.image_placeholders)
        n_bullets = len(sd.bullets)

        content_top = Inches(1.6)
        content_bottom = Inches(6.6)
        content_height = content_bottom - content_top
        left_margin = Inches(0.5)
        gap = Inches(0.3)
        total_width = Inches(11.2)

        if n_bullets == 0:
            text_width = Inches(0)
            img_width = total_width
            img_left = left_margin
        else:
            text_width = total_width * 0.42
            img_width = total_width * 0.53
            img_left = left_margin + text_width + gap

        if n_bullets > 0:
            from pptx.dml.color import RGBColor
            txBox = slide.shapes.add_textbox(left_margin, content_top, text_width, content_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_top = Pt(4)
            tf.margin_left = Pt(4)
            size_l0, size_l1, space_after, space_before = self._calc_font_size(sd.bullets)
            for i, (text, level) in enumerate(sd.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = ("  " * level + "• " if level > 0 else "• ") + text
                p.font.size = size_l0 if level == 0 else size_l1
                p.font.color.rgb = RGBColor(0x1A, 0x6B, 0xB5)
                p.space_after = space_after
                p.space_before = space_before

        if n_images > 0:
            img_gap = Inches(0.15)
            single_height = (content_height - img_gap * (n_images - 1)) / n_images
            for i, caption in enumerate(sd.image_placeholders):
                img_top = content_top + i * (single_height + img_gap)
                self._draw_image_placeholder(slide, img_left, img_top, img_width, single_height, caption)

    def _add_image_placeholders(self, slide, captions, n_bullets=0, has_code=False):
        """添加图片占位框"""
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        from lxml import etree
        from pptx.enum.shapes import MSO_SHAPE

        base_top = 1.7 + n_bullets * 0.35
        if has_code:
            base_top += 2.5
        base_top = max(base_top, 2.5)
        base_top = min(base_top, 4.8)

        page_bottom = Inches(6.6)
        n = len(captions)

        if n == 1:
            top = Inches(base_top)
            available = page_bottom - top
            height = min(Inches(2.5), available)
            if height < Inches(0.6):
                return
            self._draw_image_placeholder(slide, Inches(1.5), top, Inches(9.0), height, captions[0])
        else:
            top = Inches(base_top)
            available = page_bottom - top
            height = min(Inches(2.2), available)
            if height < Inches(0.6):
                return
            total_width = Inches(10.4)
            gap = Inches(0.2)
            w = (total_width - gap * (n - 1)) / n
            for i, caption in enumerate(captions):
                left = Inches(0.8) + i * (w + gap)
                self._draw_image_placeholder(slide, left, top, w, height, caption)

    def _draw_image_placeholder(self, slide, left, top, width, height, caption):
        """绘制单个图片占位框"""
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
        from lxml import etree
        from pptx.enum.text import PP_ALIGN

        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFF)
        box.line.color.rgb = RGBColor(0x3B, 0x8E, 0xDE)
        box.line.width = Pt(1.0)
        ln = box._element.find('.//' + qn('a:ln'))
        if ln is not None:
            prstDash = etree.SubElement(ln, qn('a:prstDash'))
            prstDash.set('val', 'dash')

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(8)
        tf.margin_left = Pt(8)

        p1 = tf.paragraphs[0]
        p1.text = "🖼"
        p1.font.size = Pt(20)
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = caption if caption else "[ 图片 ]"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0x1a, 0x6b, 0xb5)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER

    def _add_code_blocks(self, slide, code_blocks, n_bullets=0):
        """在幻灯片上添加代码块文本框"""
        from pptx.dml.color import RGBColor
        # 动态计算起始位置
        top_start = 1.8 + n_bullets * 0.35
        top_start = max(top_start, 2.5)
        top_start = min(top_start, 4.5)
        top = Inches(top_start)
        page_bottom = Inches(6.8)

        for code in code_blocks:
            lines = code.split('\n')
            n_lines = len(lines)
            height = Inches(n_lines * 0.22 + 0.3)
            max_height = page_bottom - top
            if max_height < Inches(0.5):
                break
            height = min(height, max_height)
            txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(10.4), height)
            tf = txBox.text_frame
            tf.word_wrap = True
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
            txBox.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            txBox.line.width = Pt(0.5)
            code_font_size = Pt(11) if n_lines <= 10 else Pt(10)
            for i, code_line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = code_line
                p.font.name = 'Consolas'
                p.font.size = code_font_size
                p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
                p.space_after = Pt(0)
                p.space_before = Pt(0)
            top += height + Inches(0.15)

    def _add_comparison(self, sd):
        layout = self.prs.slide_layouts[self.layouts["comparison"]]
        slide = self.prs.slides.add_slide(layout)

        # 添加标题背景色块（Comparison布局缺少这些装饰元素）
        self._add_title_background(slide)

        ph_indices = [ph.placeholder_format.idx for ph in slide.placeholders]

        if 0 in ph_indices:
            slide.placeholders[0].text = sd.title
            self._set_title_white(slide)
        if 1 in ph_indices:
            slide.placeholders[1].text = sd.left_title
        if 3 in ph_indices:
            slide.placeholders[3].text = sd.right_title

        if 2 in ph_indices:
            tf_left = slide.placeholders[2].text_frame
            tf_left.clear()
            for i, item in enumerate(sd.left_items):
                p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.space_after = Pt(8)

        if 4 in ph_indices:
            tf_right = slide.placeholders[4].text_frame
            tf_right.clear()
            for i, item in enumerate(sd.right_items):
                p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.space_after = Pt(8)

    def _add_table(self, sd):
        layout = self.prs.slide_layouts[self.layouts["table"]]
        slide = self.prs.slides.add_slide(layout)

        if 0 in [ph.placeholder_format.idx for ph in slide.placeholders]:
            slide.placeholders[0].text = sd.title

        # 如果布局缺少深蓝色标题背景，手动添加
        layout_name = layout.name.lower()
        if "spacious" in layout_name or "spacous" in layout_name:
            self._add_title_background(slide)
            self._set_title_white(slide)

        # 删除内容占位符（idx=1），避免显示 "Click to add text"
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx == 1:
                ph._element.getparent().remove(ph._element)
                break

        if sd.table_headers and sd.table_rows:
            n_cols = len(sd.table_headers)
            n_rows = len(sd.table_rows) + 1

            left = Inches(0.8)
            top = Inches(1.8)
            width = Inches(10.4)
            row_height = 0.45
            height = Inches(min(row_height * n_rows, 4.5))

            table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
            table = table_shape.table

            for i, header in enumerate(sd.table_headers):
                cell = table.cell(0, i)
                cell.text = header
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)

            for row_idx, row_data in enumerate(sd.table_rows, 1):
                for col_idx in range(min(len(row_data), n_cols)):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = row_data[col_idx]
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)

        if sd.notes:
            n_rows_total = (len(sd.table_rows) + 1) if sd.table_headers else 0
            note_top = Inches(1.8 + 0.45 * n_rows_total + 0.3)
            txBox = slide.shapes.add_textbox(Inches(0.8), note_top, Inches(10.4), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, note in enumerate(sd.notes):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = note

    def _add_section(self, sd):
        layout = self.prs.slide_layouts[self.layouts["section"]]
        slide = self.prs.slides.add_slide(layout)
        ph_indices = [ph.placeholder_format.idx for ph in slide.placeholders]
        if 0 in ph_indices:
            slide.placeholders[0].text = sd.title
        if 1 in ph_indices and sd.subtitle:
            slide.placeholders[1].text = sd.subtitle

    def save(self, output_path):
        self.prs.save(output_path)
        return output_path


# ============================================================
# Streamlit Web UI
# ============================================================

DEFAULT_TEMPLATE = "LR模板.pptx"
FORMAT_SPEC_FILE = "PPT提示词格式规范.md"
OUTPUT_DIR = "output"


def main():
    st.set_page_config(
        page_title="PPT Maker",
        page_icon="📊",
        layout="wide"
    )

    # ── 全局样式 ──────────────────────────────────────────────
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&display=swap');

    html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

    .main-header {
        background: linear-gradient(135deg, #003C71 0%, #1a6bb5 100%);
        padding: 2rem 2.5rem;
        border-radius: 12px;
        margin-bottom: 1.5rem;
        color: white;
    }
    .main-header h1 { color: white; margin: 0; font-size: 2rem; font-weight: 600; }
    .main-header p  { color: rgba(255,255,255,0.85); margin: 0.4rem 0 0; font-size: 1rem; }

    .card {
        background: #ffffff;
        border: 1px solid #e8ecf0;
        border-radius: 10px;
        padding: 1.5rem;
        margin-bottom: 1rem;
        box-shadow: 0 1px 4px rgba(0,0,0,0.06);
    }
    .card h3 { margin-top: 0; color: #003C71; font-size: 1.05rem; font-weight: 600; }

    .step-badge {
        display: inline-block;
        background: #003C71;
        color: white;
        border-radius: 50%;
        width: 26px; height: 26px;
        text-align: center; line-height: 26px;
        font-size: 0.85rem; font-weight: 600;
        margin-right: 8px;
    }

    .stButton > button {
        background: linear-gradient(135deg, #003C71, #1a6bb5);
        color: white;
        border: none;
        border-radius: 8px;
        font-weight: 600;
        font-size: 1rem;
        padding: 0.6rem 1.5rem;
        transition: opacity 0.2s;
    }
    .stButton > button:hover { opacity: 0.88; }

    .stDownloadButton > button {
        background: linear-gradient(135deg, #1a7a4a, #2ea86a);
        color: white;
        border: none;
        border-radius: 8px;
        font-weight: 600;
    }

    .success-box {
        background: #f0faf4;
        border: 1px solid #2ea86a;
        border-radius: 8px;
        padding: 1rem 1.2rem;
        color: #1a5c35;
    }
    .info-box {
        background: #f0f6ff;
        border: 1px solid #3B8EDE;
        border-radius: 8px;
        padding: 0.8rem 1.2rem;
        color: #003C71;
        font-size: 0.9rem;
    }
    </style>
    """, unsafe_allow_html=True)

    # ── 页头 ──────────────────────────────────────────────────
    st.markdown("""
    <div class="main-header">
        <h1>📊 PPT Maker</h1>
        <p>基于公司模板，从 Markdown 提示词一键生成 PowerPoint 演示文稿</p>
    </div>
    """, unsafe_allow_html=True)

    # ── 两列布局：模板 + 提示词 ───────────────────────────────
    col1, col2 = st.columns([1, 1], gap="medium")

    # ── Step 1: 模板 ──────────────────────────────────────────
    with col1:
        st.markdown('<div class="card"><h3><span class="step-badge">1</span>PPT 模板</h3>', unsafe_allow_html=True)
        template_file = st.file_uploader(
            "上传自定义模板（可选）",
            type=["pptx"],
            help="不上传则使用默认 LR 模板",
            label_visibility="collapsed"
        )
        if template_file:
            st.markdown(f'<div class="info-box">✅ 已上传：{template_file.name}</div>', unsafe_allow_html=True)
        else:
            default_exists = os.path.exists(DEFAULT_TEMPLATE)
            if default_exists:
                st.markdown('<div class="info-box">📋 使用默认 LR 模板</div>', unsafe_allow_html=True)
            else:
                st.warning("⚠️ 未找到默认模板，请上传模板文件")
        st.markdown('</div>', unsafe_allow_html=True)

    # ── Step 2: 提示词 ────────────────────────────────────────
    with col2:
        st.markdown('<div class="card"><h3><span class="step-badge">2</span>提示词内容</h3>', unsafe_allow_html=True)
        input_method = st.radio("输入方式", ["上传 .md 文件", "直接粘贴"], horizontal=True, label_visibility="collapsed")
        md_content = None

        if input_method == "上传 .md 文件":
            md_file = st.file_uploader("选择提示词文件", type=["md", "txt"], label_visibility="collapsed")
            if md_file:
                md_content = md_file.read().decode("utf-8")
                st.markdown(f'<div class="info-box">✅ 已上传：{md_file.name}</div>', unsafe_allow_html=True)
        else:
            md_content = st.text_area(
                "粘贴提示词",
                height=180,
                placeholder="## Slide 1 [cover]\n\ntitle: 演示标题\nsubtitle: 副标题\n\n## Slide 2 [content]\n\ntitle: 第一页\n\n- 要点1\n- 要点2",
                label_visibility="collapsed"
            )
        st.markdown('</div>', unsafe_allow_html=True)

    # ── 生成按钮 ──────────────────────────────────────────────
    st.markdown("<br>", unsafe_allow_html=True)
    generate_col, _ = st.columns([2, 3])
    with generate_col:
        generate_clicked = st.button("🚀 生成 PPT", type="primary", use_container_width=True)

    if generate_clicked:
        # 检查条件
        has_template = template_file is not None or os.path.exists(DEFAULT_TEMPLATE)
        if not has_template:
            st.error("❌ 请上传 PPT 模板，或将 LR模板.pptx 放在程序目录下")
            return
        if not md_content:
            st.error("❌ 请提供提示词内容")
            return

        with st.spinner("正在生成 PPT，请稍候..."):
            try:
                # 解析提示词
                parser = MarkdownParser(md_content)
                slides_data = parser.parse()
                if not slides_data:
                    st.error("❌ 提示词解析失败，未找到任何 `## Slide N` 格式的内容")
                    return

                # 准备模板
                if template_file:
                    template_file.seek(0)
                    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                        tmp.write(template_file.read())
                        template_path = tmp.name
                    cleanup_template = True
                else:
                    template_path = DEFAULT_TEMPLATE
                    cleanup_template = False

                # 生成 PPT
                gen = PPTGenerator(template_path)
                gen.generate(slides_data)

                # 保存到 output 目录，文件名带时间戳
                from datetime import datetime
                os.makedirs(OUTPUT_DIR, exist_ok=True)
                ts = datetime.now().strftime("%m%d_%H%M")
                # 从提示词第一页标题生成文件名
                base_name = slides_data[0].title[:20].strip() if slides_data else "presentation"
                # 去掉非法字符
                safe_name = re.sub(r'[\\/:*?"<>|]', '', base_name).strip() or "presentation"
                output_filename = f"{safe_name}_{ts}.pptx"
                output_path = os.path.join(OUTPUT_DIR, output_filename)
                gen.save(output_path)

                # 读取文件供下载
                with open(output_path, "rb") as f:
                    pptx_bytes = f.read()

                if cleanup_template:
                    os.unlink(template_path)

                # 成功提示
                st.markdown(f"""
                <div class="success-box">
                    ✅ <strong>生成成功！</strong>共 {len(slides_data)} 页<br>
                    📁 已保存至 <code>output/{output_filename}</code>
                </div>
                """, unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)

                # 下载按钮
                st.download_button(
                    label=f"📥 下载 {output_filename}",
                    data=pptx_bytes,
                    file_name=output_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )

            except Exception as e:
                st.error(f"❌ 生成失败：{str(e)}")
                st.exception(e)

    # ── 底部格式帮助 ──────────────────────────────────────────
    st.markdown("<br>", unsafe_allow_html=True)

    # ── Step 3: 参考资源 ──────────────────────────────────────
    st.markdown('<div class="card"><h3><span class="step-badge">3</span>参考资源</h3>', unsafe_allow_html=True)
    res_col1, res_col2 = st.columns([1, 2], gap="medium")
    with res_col1:
        if os.path.exists(FORMAT_SPEC_FILE):
            with open(FORMAT_SPEC_FILE, "rb") as f:
                st.download_button(
                    label="📄 下载提示词格式规范",
                    data=f.read(),
                    file_name="PPT提示词格式规范.md",
                    mime="text/markdown",
                    use_container_width=True
                )
        else:
            st.info("格式规范文件未找到")
    with res_col2:
        if md_content:
            parser = MarkdownParser(md_content)
            slides = parser.parse()
            if slides:
                for sd in slides:
                    icon = {"cover": "🎯", "content": "📝", "table": "📊",
                            "comparison": "⚖️", "section": "📌", "end": "🏁"}.get(sd.layout, "📄")
                    st.markdown(f"{icon} **{sd.index}** {sd.title[:40] or '(无标题)'}")
            else:
                st.warning("未解析到幻灯片，请检查格式")
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    with st.expander("❓ 提示词格式速查"):
        st.markdown("""
| 布局标记 | 用途 | 必填字段 |
|----------|------|---------|
| `[cover]` | 封面页 | `title:`, `subtitle:` |
| `[content]` | 正文页 | `title:`, `- 要点` |
| `[comparison]` | 双栏对比 | `title:`, `left_title:`, `left:`, `right_title:`, `right:` |
| `[table]` | 表格页 | `title:`, Markdown 表格 |
| `[section]` | 章节页 | `title:`, `subtitle:` |
| `[end]` | 结束页 | `title:` |

**代码块**：在正文页中用 ` ``` ` 包裹代码，自动渲染为等宽字体文本框。

**图片占位**：用 `image: 说明文字` 预留图片位置，后期手工替换。

**自动字体缩放**：内容越多字号越小，保证不溢出页面。

**多格式支持**：`## Slide 1 [cover]` / `## Slide 1 — Cover` / `## 第1页 封面` 均可识别。
        """)


if __name__ == "__main__":
    main()

